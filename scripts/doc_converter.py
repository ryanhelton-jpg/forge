#!/usr/bin/env python3
"""
Multi-Format Document Converter
Usage: python3 doc_converter.py <format> <input_file> [--title "Title"] [--template professional]

Formats: docx, xlsx, pptx, pdf
"""

import argparse
import json
import sys
from pathlib import Path
from typing import Dict, Any, Optional

class DocumentConverter:
    """Professional document converter supporting multiple output formats."""
    
    TEMPLATES = {
        'professional': {'colors': ['#1f4e79', '#2e75b6'], 'font': 'Calibri'},
        'academic': {'colors': ['#2f3e46', '#354f52'], 'font': 'Times New Roman'},
        'executive': {'colors': ['#2d3436', '#74b9ff'], 'font': 'Arial'},
        'casual': {'colors': ['#6c5ce7', '#a29bfe'], 'font': 'Arial'},
    }

    def __init__(self, template: str = 'professional'):
        self.template = self.TEMPLATES.get(template, self.TEMPLATES['professional'])

    def to_docx(self, content: str, output_path: str, title: str = "Document", 
                metadata: Optional[Dict] = None) -> str:
        """Convert content to Word document."""
        try:
            from docx import Document
            from docx.shared import Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
        except ImportError:
            return "❌ python-docx not installed. Run: pip install python-docx"

        try:
            doc = Document()
            
            # Set metadata
            if metadata:
                doc.core_properties.author = metadata.get('author', 'Generated')
                doc.core_properties.title = title
                doc.core_properties.subject = metadata.get('subject', '')
            
            # Add title
            title_para = doc.add_heading(title, 0)
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Process markdown content
            lines = content.split('\n')
            for line in lines:
                line = line.rstrip()
                if line.startswith('# '):
                    doc.add_heading(line[2:], level=1)
                elif line.startswith('## '):
                    doc.add_heading(line[3:], level=2)
                elif line.startswith('### '):
                    doc.add_heading(line[4:], level=3)
                elif line.startswith('- ') or line.startswith('* '):
                    doc.add_paragraph(line[2:], style='List Bullet')
                elif line.strip():
                    doc.add_paragraph(line)
            
            doc.save(output_path)
            return f"✅ Word document saved: {output_path}"
        except Exception as e:
            return f"❌ Error creating Word document: {str(e)}"

    def to_xlsx(self, content: str, output_path: str, title: str = "Spreadsheet") -> str:
        """Convert content to Excel spreadsheet."""
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment
        except ImportError:
            return "❌ openpyxl not installed. Run: pip install openpyxl"

        try:
            wb = Workbook()
            ws = wb.active
            ws.title = title[:31]  # Excel sheet name limit
            
            # Parse content as CSV-like or table data
            lines = content.strip().split('\n')
            for row_idx, line in enumerate(lines, 1):
                # Try to split by common delimiters
                if '\t' in line:
                    cells = line.split('\t')
                elif '|' in line:
                    cells = [c.strip() for c in line.split('|') if c.strip()]
                elif ',' in line:
                    cells = line.split(',')
                else:
                    cells = [line]
                
                for col_idx, cell in enumerate(cells, 1):
                    ws.cell(row=row_idx, column=col_idx, value=cell.strip())
                    
                    # Style header row
                    if row_idx == 1:
                        ws.cell(row=row_idx, column=col_idx).font = Font(bold=True)
            
            # Auto-size columns
            for column in ws.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                ws.column_dimensions[column_letter].width = min(max_length + 2, 50)
            
            wb.save(output_path)
            return f"✅ Excel spreadsheet saved: {output_path}"
        except Exception as e:
            return f"❌ Error creating Excel file: {str(e)}"

    def to_pptx(self, content: str, output_path: str, title: str = "Presentation") -> str:
        """Convert content to PowerPoint presentation."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            return "❌ python-pptx not installed. Run: pip install python-pptx"

        try:
            prs = Presentation()
            
            # Title slide
            title_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = title
            
            # Parse content into slides (split on ## headers)
            sections = content.split('\n## ')
            
            for i, section in enumerate(sections):
                if i == 0 and not section.startswith('## '):
                    # First section might be intro text
                    if section.strip() and not section.startswith('# '):
                        bullet_layout = prs.slide_layouts[1]
                        slide = prs.slides.add_slide(bullet_layout)
                        slide.shapes.title.text = "Overview"
                        body = slide.shapes.placeholders[1]
                        tf = body.text_frame
                        tf.text = section.strip()[:500]
                    continue
                
                lines = section.split('\n')
                slide_title = lines[0].replace('## ', '').strip()
                slide_content = '\n'.join(lines[1:]).strip()
                
                bullet_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_layout)
                slide.shapes.title.text = slide_title
                
                if slide_content:
                    body = slide.shapes.placeholders[1]
                    tf = body.text_frame
                    tf.text = slide_content[:500]
            
            prs.save(output_path)
            return f"✅ PowerPoint saved: {output_path}"
        except Exception as e:
            return f"❌ Error creating PowerPoint: {str(e)}"

    def to_pdf(self, content: str, output_path: str, title: str = "Document") -> str:
        """Convert content to PDF via markdown → HTML → PDF."""
        try:
            import subprocess
            import tempfile
            import markdown
            
            # Convert markdown to HTML
            html_content = f"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{title}</title>
    <style>
        body {{ font-family: {self.template['font']}, sans-serif; margin: 40px; line-height: 1.6; }}
        h1, h2, h3 {{ color: {self.template['colors'][0]}; }}
        h1 {{ border-bottom: 2px solid {self.template['colors'][1]}; padding-bottom: 10px; }}
        code {{ background: #f4f4f4; padding: 2px 6px; border-radius: 3px; }}
        pre {{ background: #f4f4f4; padding: 15px; border-radius: 5px; overflow-x: auto; }}
    </style>
</head>
<body>
    <h1>{title}</h1>
    {markdown.markdown(content, extensions=['tables', 'fenced_code'])}
</body>
</html>
"""
            # Write temp HTML
            with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False) as f:
                f.write(html_content)
                html_path = f.name
            
            # Try wkhtmltopdf first, then weasyprint
            try:
                result = subprocess.run(
                    ['wkhtmltopdf', '--quiet', html_path, output_path],
                    capture_output=True, text=True, timeout=30
                )
                if result.returncode == 0:
                    Path(html_path).unlink()
                    return f"✅ PDF saved: {output_path}"
            except FileNotFoundError:
                pass
            
            # Try weasyprint as fallback
            try:
                from weasyprint import HTML
                HTML(filename=html_path).write_pdf(output_path)
                Path(html_path).unlink()
                return f"✅ PDF saved: {output_path}"
            except ImportError:
                pass
            
            # Try pandoc as last resort
            try:
                result = subprocess.run(
                    ['pandoc', html_path, '-o', output_path],
                    capture_output=True, text=True, timeout=30
                )
                Path(html_path).unlink()
                if result.returncode == 0:
                    return f"✅ PDF saved: {output_path}"
                return f"❌ Pandoc error: {result.stderr}"
            except FileNotFoundError:
                Path(html_path).unlink()
                return "❌ No PDF converter found. Install wkhtmltopdf, weasyprint, or pandoc."
                
        except ImportError:
            return "❌ markdown not installed. Run: pip install markdown"
        except Exception as e:
            return f"❌ Error creating PDF: {str(e)}"


def main():
    parser = argparse.ArgumentParser(description='Convert documents to various formats')
    parser.add_argument('format', choices=['docx', 'xlsx', 'pptx', 'pdf'], 
                        help='Output format')
    parser.add_argument('input_file', help='Input file (text/markdown)')
    parser.add_argument('-o', '--output', help='Output file path')
    parser.add_argument('-t', '--title', default='Document', help='Document title')
    parser.add_argument('--template', default='professional',
                        choices=['professional', 'academic', 'executive', 'casual'],
                        help='Template style')
    parser.add_argument('--metadata', help='JSON metadata string')
    
    args = parser.parse_args()
    
    # Read input
    input_path = Path(args.input_file)
    if not input_path.exists():
        print(f"❌ Input file not found: {args.input_file}")
        sys.exit(1)
    
    content = input_path.read_text()
    
    # Determine output path
    if args.output:
        output_path = args.output
    else:
        output_path = str(input_path.with_suffix(f'.{args.format}'))
    
    # Parse metadata if provided
    metadata = None
    if args.metadata:
        try:
            metadata = json.loads(args.metadata)
        except json.JSONDecodeError:
            print("❌ Invalid metadata JSON")
            sys.exit(1)
    
    # Convert
    converter = DocumentConverter(template=args.template)
    
    if args.format == 'docx':
        result = converter.to_docx(content, output_path, args.title, metadata)
    elif args.format == 'xlsx':
        result = converter.to_xlsx(content, output_path, args.title)
    elif args.format == 'pptx':
        result = converter.to_pptx(content, output_path, args.title)
    elif args.format == 'pdf':
        result = converter.to_pdf(content, output_path, args.title)
    
    print(result)
    sys.exit(0 if result.startswith('✅') else 1)


if __name__ == '__main__':
    main()
